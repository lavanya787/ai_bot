import os
import re
import pickle
import faiss
import numpy as np
from typing import List
from sentence_transformers import SentenceTransformer
import fitz  # PyMuPDF
import docx
import speech_recognition as sr
from PIL import Image
import pytesseract
import logging
from concurrent.futures import ThreadPoolExecutor
import pptx
from pptx import Presentation

try:
    import spacy
    nlp = spacy.load("en_core_web_sm")
    SPACY_AVAILABLE = True
except Exception:
    nlp = None
    SPACY_AVAILABLE = False

logger = logging.getLogger("processor")
logging.basicConfig(level=logging.INFO)

class DocumentProcessor:
    def __init__(self, base_dir="faiss_logs", model_name="all-MiniLM-L6-v2", device="cpu"):
        self.base_dir = base_dir
        os.makedirs(base_dir, exist_ok=True)
        self.model = SentenceTransformer(model_name)

    def file_type_detector(self, file_name: str) -> str:
        ext = os.path.splitext(file_name)[-1].lower()
        return {
            ".txt": "text",
            ".pdf": "pdf",
            ".docx": "docx",
            ".wav": "audio",
            ".mp3": "audio",
            ".m4a": "audio",
            ".png": "image",
            ".jpg": "image",
            ".jpeg": "image",
            ".csv": "csv",
            ".xls": "excel",
            ".xlsx": "excel",
            ".json": "json",
            ".ppt": "ppt",
            ".pptx": "ppt"
        }.get(ext, "unknown")

    def read_file(self, file, file_type: str) -> str:
        file.seek(0)
        if file_type == "text":
            return file.read().decode("utf-8", errors="ignore")
        elif file_type == "pdf":
            text = ""
            with fitz.open(stream=file.read(), filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text()
            return text
        elif file_type == "docx":
            doc = docx.Document(file)
            return "\n".join([para.text for para in doc.paragraphs])
        elif file_type == "audio":
            recognizer = sr.Recognizer()
            with sr.AudioFile(file) as source:
                audio_data = recognizer.record(source)
                try:
                    return recognizer.recognize_google(audio_data)
                except sr.UnknownValueError:
                    return ""
        elif file_type == "image":
            image = Image.open(file)
            return pytesseract.image_to_string(image)
        elif file_type == "csv":
            import pandas as pd
            df = pd.read_csv(file)
            return "\n".join(df.astype(str).apply(lambda row: " | ".join(row), axis=1))
        elif file_type == "excel":
            import pandas as pd
            df = pd.read_excel(file)
            return "\n".join(df.astype(str).apply(lambda row: " | ".join(row), axis=1))
        elif file_type == "json":
            import json
            data = json.load(file)
            return json.dumps(data, indent=2)
        elif file_type == "ppt":
            prs = Presentation(file)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides.append(shape.text.strip())
            return "\n".join(slides)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def chunk_text(self, text: str, chunk_size: int = 1200) -> List[str]:
        if not text or len(text.strip()) < 50:
            return []

        try:
            if SPACY_AVAILABLE and nlp:
                doc = nlp(text)
                sentences = [sent.text.strip() for sent in doc.sents]
            else:
                sentences = re.split(r'(?<=[.!?])\s+', text.strip())
        except Exception as e:
            logger.warning(f"Chunk fallback due to error: {e}")
            sentences = text.split(". ")

        chunks = []
        with ThreadPoolExecutor(max_workers=4) as executor:
            futures = []
            i = 0
            while i < len(sentences):
                chunk = []
                total_len = 0
                while i < len(sentences) and total_len + len(sentences[i]) < chunk_size:
                    chunk.append(sentences[i])
                    total_len += len(sentences[i])
                    i += 1
                current_chunk = chunk.copy()
                futures.append(executor.submit(lambda x: " ".join(x).strip(), current_chunk))

            for future in futures:
                result = future.result()
                if len(result.split()) >= 10:
                    chunks.append(result)

        return chunks

    def chunk_by_heading(self, text: str) -> List[str]:
        pattern = r'\n\s*(?:Chapter|CHAPTER|\d+\.\s+)[^\n]+\n'
        matches = list(re.finditer(pattern, text))
        if not matches:
            return self.chunk_text(text)

        sections = []
        for i, match in enumerate(matches):
            start = match.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            section = text[start:end].strip()
            if section:
                title = match.group(0).strip()
                sections.append(f"{title}\n{section}")
        return sections

    def build_faiss_index(self, chunks: List[str], file_name: str, batch_size=32) -> str:
        file_id = os.path.splitext(file_name)[0].replace(" ", "_").lower()
        folder = os.path.join(self.base_dir, file_id)

        if os.path.exists(os.path.join(folder, "faiss.index")):
            return folder

        os.makedirs(folder, exist_ok=True)

        embeddings = self.model.encode(
            chunks,
            batch_size=batch_size,
            show_progress_bar=True,
            convert_to_numpy=True
        )
        dim = embeddings.shape[1]
        index = faiss.IndexFlatL2(dim)
        index.add(np.array(embeddings).astype("float32"))

        faiss.write_index(index, os.path.join(folder, "faiss.index"))
        with open(os.path.join(folder, "chunks.pkl"), "wb") as f:
            pickle.dump(chunks, f)

        return folder

    def process_and_index_file(self, file, file_name: str, llm_handler=None) -> str:
        file_type = self.file_type_detector(file_name)
        if file_type == "unknown":
            raise ValueError(f"Unsupported file format: {file_name}")

        text = self.read_file(file, file_type)
        chunks = self.chunk_by_heading(text) if "chapter" in text.lower() else self.chunk_text(text)
        if not chunks:
            raise ValueError(f"No valid content extracted from file: {file_name}")

        if llm_handler:
            for i, chunk in enumerate(chunks):
                llm_handler.index_document(f"{file_name}_chunk_{i}", chunk)

        return self.build_faiss_index(chunks, file_name)

    def search(self, query: str, file_name: str, top_k: int = 5) -> List[str]:
        file_id = os.path.splitext(file_name)[0].replace(" ", "_").lower()
        folder = os.path.join(self.base_dir, file_id)
        if not os.path.exists(folder):
            raise FileNotFoundError(f"No index found for {file_id}")

        index = faiss.read_index(os.path.join(folder, "faiss.index"))
        with open(os.path.join(folder, "chunks.pkl"), "rb") as f:
            chunks = pickle.load(f)

        query_embedding = self.model.encode([query], convert_to_numpy=True)
        D, I = index.search(np.array(query_embedding).astype("float32"), top_k)

        return [chunks[i] for i in I[0] if i < len(chunks)]
